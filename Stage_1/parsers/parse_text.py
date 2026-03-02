import logging
import re
from pathlib import Path
from Stage_1.ParseResult import ParseResult
import Stage_1.registry as registry

logger = logging.getLogger(__name__)

# Returns standardized UTF-8 string

"""
Text parsers.

Handles: plain text, code, markup, documents (PDF, DOCX, PPTX).
Each function returns ParseResult(modality="text", text=...).

Multi-modal detection:
    PDF, DOCX, and PPTX parsers check for embedded images/tables while
    they're already inside the file. This is nearly free since we're
    iterating the file structure anyway. Detected extras go into
    result.also_contains so the orchestrator can queue follow-up work.
"""


# ===================================================================
# CONFIG
# ===================================================================

DEFAULT_MAX_CHARS = 500_000  # ~125k tokens


def _max_chars(config: dict) -> int:
    return config.get("max_chars", DEFAULT_MAX_CHARS)


def _clean_text(text: str) -> str:
    """Normalize whitespace and remove junk."""
    if not text:
        return ""
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ===================================================================
# PLAIN TEXT / CODE
# ===================================================================

def parse_plaintext(path: str, config: dict) -> ParseResult:
    """Read any UTF-8 text file. Falls back to latin-1."""
    try:
        limit = _max_chars(config)
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read(limit)
        except UnicodeDecodeError:
            with open(path, "r", encoding="latin-1") as f:
                content = f.read(limit)

        content = _clean_text(content)

        return ParseResult(
            modality="text",
            output=content,
            metadata={"char_count": len(content)},
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="text")


registry.register([
    ".txt", ".md", ".markdown", ".rst", ".tex", ".log", ".rtf",
    ".json", ".yaml", ".yml", ".xml",
    ".ini", ".toml", ".cfg", ".env",
    ".py", ".js", ".jsx", ".ts", ".tsx",
    ".html", ".htm", ".css", ".scss",
    ".c", ".cpp", ".h", ".hpp",
    ".java", ".cs", ".php", ".rb",
    ".go", ".rs", ".swift", ".kt",
    ".sql", ".sh", ".bat", ".ps1",
    ".r", ".m", ".scala", ".lua",
], "text", parse_plaintext)


# ===================================================================
# PDF
# ===================================================================

def parse_pdf_text(path: str, config: dict) -> ParseResult:
    """
    Extract text from a PDF.

    Also detects embedded images and scanned pages (no extractable text).
    These go into also_contains so the orchestrator can queue image
    extraction or OCR tasks.
    """
    try:
        from pdfminer.high_level import extract_text
    except ImportError:
        logger.debug("pdfminer.six not installed")
        return ParseResult.failed("pdfminer.six not installed", modality="text")

    try:
        limit = _max_chars(config)
        text = extract_text(path)
        text = _clean_text(text[:limit])

        # --- Cheap multi-modal detection ---
        metadata = {"char_count": len(text)}
        also_contains = []
        image_count = 0
        page_count = 0

        try:
            from pdfminer.pdfpage import PDFPage
            from pdfminer.pdfparser import PDFParser
            from pdfminer.pdfdocument import PDFDocument
            from pdfminer.psparser import LIT

            with open(path, "rb") as f:
                parser = PDFParser(f)
                doc = PDFDocument(parser)

                for page in PDFPage.create_pages(doc):
                    page_count += 1
                    resources = page.resources
                    if not resources:
                        continue
                    xobjects = resources.get("XObject")
                    if not xobjects:
                        continue
                    # Resolve indirect references
                    if hasattr(xobjects, "resolve"):
                        xobjects = xobjects.resolve()
                    if isinstance(xobjects, dict):
                        for _, obj in xobjects.items():
                            if hasattr(obj, "resolve"):
                                obj = obj.resolve()
                            if isinstance(obj, dict):
                                subtype = obj.get("Subtype")
                                if subtype is LIT("Image") or subtype == "/Image":
                                    image_count += 1
        except Exception as e:
            # Detection failed — not critical, just skip
            logger.debug(f"PDF multi-modal detection failed: {e}")

        metadata["page_count"] = page_count
        metadata["image_count"] = image_count
        metadata["has_images"] = image_count > 0
        metadata["is_scanned"] = len(text.strip()) < 50 and page_count > 0

        if image_count > 0:
            also_contains.append("image")
        if metadata.get("is_scanned"):
            # Scanned PDFs need OCR — flag as needing image processing
            if "image" not in also_contains:
                also_contains.append("image")

        return ParseResult(
            modality="text",
            output=text,
            metadata=metadata,
            also_contains=also_contains,
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="text")


registry.register(".pdf", "text", parse_pdf_text)


def parse_pdf_image(path: str, config: dict) -> ParseResult:
    """
    Extract embedded images from a PDF as PIL.Image objects.

    This is called when a task requests parse("report.pdf", "image").
    Returns a list of PIL images — one per embedded image found.
    """
    try:
        from pdfminer.pdfpage import PDFPage
        from pdfminer.pdfparser import PDFParser
        from pdfminer.pdfdocument import PDFDocument
        from PIL import Image
        import io
    except ImportError as e:
        logger.debug(f"Missing dependency: {e}")
        return ParseResult.failed(f"Missing dependency: {e}", modality="image")

    try:
        images = []
        max_images = config.get("max_images", 50)

        with open(path, "rb") as f:
            parser = PDFParser(f)
            doc = PDFDocument(parser)

            for page in PDFPage.create_pages(doc):
                if len(images) >= max_images:
                    break
                resources = page.resources
                if not resources:
                    continue
                xobjects = resources.get("XObject")
                if not xobjects:
                    continue
                if hasattr(xobjects, "resolve"):
                    xobjects = xobjects.resolve()
                if not isinstance(xobjects, dict):
                    continue
                for _, obj in xobjects.items():
                    if len(images) >= max_images:
                        break
                    if hasattr(obj, "resolve"):
                        obj = obj.resolve()
                    if not isinstance(obj, dict):
                        continue
                    # Try to extract image data
                    try:
                        data = obj.get_data()
                        img = Image.open(io.BytesIO(data))
                        images.append(img)
                    except Exception:
                        # Not all XObjects are extractable this way
                        continue

        if not images:
            return ParseResult.failed(
                "No extractable images found in PDF", modality="image"
            )

        return ParseResult(
            modality="image",
            output=images,
            metadata={"image_count": len(images), "source_format": "pdf"},
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="image")


registry.register(".pdf", "image", parse_pdf_image)


# ===================================================================
# DOCX
# ===================================================================

def parse_docx_text(path: str, config: dict) -> ParseResult:
    """Extract text from a Word document. Detects embedded images."""
    try:
        from docx import Document
    except ImportError:
        logger.debug("python-docx not installed")
        return ParseResult.failed("python-docx not installed", modality="text")

    try:
        limit = _max_chars(config)
        doc = Document(path)

        paragraphs = []
        current_len = 0
        for para in doc.paragraphs:
            paragraphs.append(para.text)
            current_len += len(para.text)
            if current_len > limit:
                break

        content = _clean_text("\n".join(paragraphs)[:limit])

        # --- Multi-modal detection ---
        also_contains = []
        image_count = 0
        has_tables = len(doc.tables) > 0

        # Check for image relationships
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_count += 1

        metadata = {
            "char_count": len(content),
            "paragraph_count": len(doc.paragraphs),
            "image_count": image_count,
            "has_images": image_count > 0,
            "has_tables": has_tables,
            "table_count": len(doc.tables),
        }

        if image_count > 0:
            also_contains.append("image")
        if has_tables:
            also_contains.append("tabular")

        return ParseResult(
            modality="text",
            output=content,
            metadata=metadata,
            also_contains=also_contains,
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="text")


registry.register([".docx", ".doc"], "text", parse_docx_text)


def parse_docx_image(path: str, config: dict) -> ParseResult:
    """Extract embedded images from a DOCX as PIL.Image objects."""
    try:
        from docx import Document
        from PIL import Image
        import io
    except ImportError as e:
        logger.debug(f"Missing dependency: {e}")
        return ParseResult.failed(f"Missing dependency: {e}", modality="image")

    try:
        doc = Document(path)
        images = []
        max_images = config.get("max_images", 50)

        for rel in doc.part.rels.values():
            if len(images) >= max_images:
                break
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    img = Image.open(io.BytesIO(image_data))
                    images.append(img)
                except Exception:
                    continue

        if not images:
            return ParseResult.failed(
                "No extractable images found in DOCX", modality="image"
            )

        return ParseResult(
            modality="image",
            output=images,
            metadata={"image_count": len(images), "source_format": "docx"},
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="image")


registry.register([".docx", ".doc"], "image", parse_docx_image)


# ===================================================================
# PPTX
# ===================================================================

def parse_pptx_text(path: str, config: dict) -> ParseResult:
    """Extract text from a PowerPoint. Detects embedded images."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.debug("python-pptx not installed")
        return ParseResult.failed("python-pptx not installed", modality="text")

    try:
        limit = _max_chars(config)
        prs = Presentation(path)

        text_runs = []
        current_len = 0
        image_count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text)
                    current_len += len(shape.text)
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image_count += 1
            if current_len > limit:
                break

        content = _clean_text("\n".join(text_runs)[:limit])

        also_contains = []
        if image_count > 0:
            also_contains.append("image")

        return ParseResult(
            modality="text",
            output=content,
            metadata={
                "char_count": len(content),
                "slide_count": len(prs.slides),
                "image_count": image_count,
                "has_images": image_count > 0,
            },
            also_contains=also_contains,
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="text")


registry.register(".pptx", "text", parse_pptx_text)


def parse_pptx_image(path: str, config: dict) -> ParseResult:
    """Extract embedded images from a PPTX as PIL.Image objects."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image
        import io
    except ImportError as e:
        logger.debug(f"Missing dependency: {e}")
        return ParseResult.failed(f"Missing dependency: {e}", modality="image")

    try:
        prs = Presentation(path)
        images = []
        max_images = config.get("max_images", 50)

        for slide in prs.slides:
            if len(images) >= max_images:
                break
            for shape in slide.shapes:
                if len(images) >= max_images:
                    break
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_data = shape.image.blob
                        img = Image.open(io.BytesIO(image_data))
                        images.append(img)
                    except Exception:
                        continue

        if not images:
            return ParseResult.failed(
                "No extractable images found in PPTX", modality="image"
            )

        return ParseResult(
            modality="image",
            output=images,
            metadata={"image_count": len(images), "source_format": "pptx"},
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="image")


registry.register(".pptx", "image", parse_pptx_image)