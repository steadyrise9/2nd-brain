import logging
import re
import time
from pathlib import Path
from Stage_1.services.ParseResult import ParseResult
from Stage_1.services import parser_registry as registry

logger = logging.getLogger("ParseText")

# Returns standardized UTF-8 string

"""
Text parsers.

Handles: plain text, code, markup, documents (PDF, DOCX, PPTX).
Each function returns ParseResult(modality="text", text=...).

Multi-modal detection:
    PDF, DOCX, and PPTX parsers check for embedded images/tables while
    they're already inside the file. This is nearly free since we're
    iterating the file structure anyway. Detected extras go into
    result.also_contains so the orchestrator can queue follow-up work.
"""


# ===================================================================
# CONFIG
# ===================================================================

DEFAULT_MAX_CHARS = 500_000  # ~125k tokens


def _max_chars(config: dict) -> int:
    return config.get("max_chars", DEFAULT_MAX_CHARS)


def _clean_text(text: str, preserve_indent: bool = False) -> str:
    """Normalize whitespace and remove junk.

    If preserve_indent is True, only collapse horizontal whitespace
    within lines (not leading whitespace), keeping indentation intact.
    """
    if not text:
        return ""
    if preserve_indent:
        # Collapse runs of spaces/tabs mid-line only, keep leading whitespace
        text = re.sub(r"(?<=\S)[ \t]+", " ", text)
    else:
        text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ===================================================================
# PLAIN TEXT / CODE
# ===================================================================

def parse_plaintext(path: str, config: dict, services: dict = None) -> ParseResult:
    """Read any UTF-8 text file. Falls back to latin-1."""
    try:
        limit = _max_chars(config)
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read(limit)
        except UnicodeDecodeError:
            with open(path, "r", encoding="latin-1") as f:
                content = f.read(limit)

        # Preserve indentation for code/config files
        _CODE_SUFFIXES = {
            ".py", ".js", ".jsx", ".ts", ".tsx", ".html", ".htm", ".css", ".scss",
            ".c", ".cpp", ".h", ".hpp", ".java", ".cs", ".php", ".rb",
            ".go", ".rs", ".swift", ".kt", ".sql", ".sh", ".bat", ".ps1",
            ".r", ".m", ".scala", ".lua", ".json", ".yaml", ".yml", ".xml",
            ".ini", ".toml", ".cfg", ".env", ".log",
        }
        is_code = Path(path).suffix.lower() in _CODE_SUFFIXES
        content = _clean_text(content, preserve_indent=is_code)

        return ParseResult(
            modality="text",
            output=content,
            metadata={"char_count": len(content)},
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="text")


registry.register([
    ".txt", ".md", ".markdown", ".rst", ".tex", ".log", ".rtf",
    ".json", ".yaml", ".yml", ".xml",
    ".ini", ".toml", ".cfg", ".env",
    ".py", ".js", ".jsx", ".ts", ".tsx",
    ".html", ".htm", ".css", ".scss",
    ".c", ".cpp", ".h", ".hpp",
    ".java", ".cs", ".php", ".rb",
    ".go", ".rs", ".swift", ".kt",
    ".sql", ".sh", ".bat", ".ps1",
    ".r", ".m", ".scala", ".lua",
], "text", parse_plaintext)


# ===================================================================
# PDF
# ===================================================================

def parse_pdf_text(path: str, config: dict, services: dict = None) -> ParseResult:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.debug("PyMuPDF not installed")
        return ParseResult.failed("PyMuPDF not installed", modality="text")

    try:
        t0 = time.time()
        limit = _max_chars(config)

        with fitz.open(path) as doc:
            # Extract text
            text_parts = []
            current_len = 0
            image_count = 0
            has_tables = False

            for page in doc:
                # Text extraction
                page_text = page.get_text()
                text_parts.append(page_text)
                current_len += len(page_text)

                # Image detection (cheap — reads page structure, not pixel data)
                image_count += len(page.get_images(full=False))

                # Table detection
                if not has_tables:
                    tables = page.find_tables()
                    if tables.tables:
                        has_tables = True

                if current_len > limit:
                    break

            text = _clean_text("".join(text_parts)[:limit])

            also_contains = []
            if image_count > 0:
                also_contains.append("image")
            if has_tables:
                also_contains.append("tabular")

            # A PDF with almost no text but multiple pages is likely a scanned document.
            # Flag it so the orchestrator queues OCR via the image modality.
            is_scanned = len(text.strip()) < 50 and len(doc) > 0

            if is_scanned and "image" not in also_contains:
                also_contains.append("image")

            metadata = {
                "char_count": len(text),
                "page_count": len(doc),
                "image_count": image_count,
                "has_images": image_count > 0,
                "has_tables": has_tables,
                "is_scanned": is_scanned,
            }

        logger.debug(
            f"PDF parsed: {Path(path).name} — {metadata['page_count']} pages, "
            f"{len(text)} chars in {time.time() - t0:.2f}s"
        )
        return ParseResult(
            modality="text",
            output=text,
            metadata=metadata,
            also_contains=also_contains,
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="text")


registry.register(".pdf", "text", parse_pdf_text)


def parse_pdf_image(path: str, config: dict, services: dict = None) -> ParseResult:
    try:
        import fitz
        from PIL import Image
        import io
    except ImportError as e:
        logger.debug(f"Missing dependency: {e}")
        return ParseResult.failed(f"Missing dependency: {e}", modality="image")

    try:
        with fitz.open(path) as doc:
            images = []
            max_images = config.get("max_images", 50)

            for page in doc:
                if len(images) >= max_images:
                    break
                for img_info in page.get_images(full=True):
                    if len(images) >= max_images:
                        break
                    xref = img_info[0]
                    try:
                        base_image = doc.extract_image(xref)
                        img = Image.open(io.BytesIO(base_image["image"]))
                        images.append(img)
                    except Exception:
                        continue

        if not images:
            return ParseResult.failed(
                "No extractable images found in PDF", modality="image"
            )

        return ParseResult(
            modality="image",
            output=images,
            metadata={"image_count": len(images), "source_format": "pdf"},
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="image")


registry.register(".pdf", "image", parse_pdf_image)


def parse_pdf_tables(path: str, config: dict, services: dict = None) -> ParseResult:
    """Extract tables from a PDF as DataFrames."""
    try:
        import fitz
        import pandas as pd
    except ImportError as e:
        logger.debug(f"Missing dependency: {e}")
        return ParseResult.failed(f"Missing dependency: {e}", modality="tabular")

    try:
        with fitz.open(path) as doc:
            all_tables = {}
            table_idx = 0

            for page_num, page in enumerate(doc):
                found = page.find_tables()
                for table in found.tables:
                    df = table.to_pandas()
                    # Skip empty tables
                    if df.empty:
                        continue
                    key = f"page_{page_num + 1}_table_{table_idx + 1}"
                    all_tables[key] = df
                    table_idx += 1

        if not all_tables:
            return ParseResult.failed(
                "No extractable tables found in PDF", modality="tabular"
            )

        table_meta = {}
        total_rows = 0
        for name, df in all_tables.items():
            total_rows += len(df)
            table_meta[name] = {
                "row_count": len(df),
                "column_count": len(df.columns),
                "columns": list(df.columns),
                "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
            }

        return ParseResult(
            modality="tabular",
            output=all_tables,
            metadata={
                "total_rows": total_rows,
                "table_count": len(all_tables),
                "table_names": list(all_tables.keys()),
                "tables": table_meta,
                "source_format": "pdf",
            },
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="tabular")


registry.register(".pdf", "tabular", parse_pdf_tables)

# ===================================================================
# DOCX
# ===================================================================

def parse_docx_text(path: str, config: dict, services: dict = None) -> ParseResult:
    """Extract text from a Word document. Detects embedded images."""
    try:
        from docx import Document
    except ImportError:
        logger.debug("python-docx not installed")
        return ParseResult.failed("python-docx not installed", modality="text")

    try:
        t0 = time.time()
        limit = _max_chars(config)
        doc = Document(path)

        paragraphs = []
        current_len = 0
        for para in doc.paragraphs:
            paragraphs.append(para.text)
            current_len += len(para.text)
            if current_len > limit:
                break

        content = _clean_text("\n".join(paragraphs)[:limit])

        # --- Multi-modal detection ---
        also_contains = []
        image_count = 0
        has_tables = len(doc.tables) > 0

        # Check for image relationships
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_count += 1

        metadata = {
            "char_count": len(content),
            "paragraph_count": len(doc.paragraphs),
            "image_count": image_count,
            "has_images": image_count > 0,
            "has_tables": has_tables,
            "table_count": len(doc.tables),
        }

        if image_count > 0:
            also_contains.append("image")
        if has_tables:
            also_contains.append("tabular")

        logger.debug(
            f"DOCX parsed: {Path(path).name} — {len(doc.paragraphs)} paragraphs, "
            f"{len(content)} chars in {time.time() - t0:.2f}s"
        )
        return ParseResult(
            modality="text",
            output=content,
            metadata=metadata,
            also_contains=also_contains,
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="text")


registry.register([".docx", ".doc"], "text", parse_docx_text)


def parse_docx_image(path: str, config: dict, services: dict = None) -> ParseResult:
    """Extract embedded images from a DOCX as PIL.Image objects."""
    try:
        from docx import Document
        from PIL import Image
        import io
    except ImportError as e:
        logger.debug(f"Missing dependency: {e}")
        return ParseResult.failed(f"Missing dependency: {e}", modality="image")

    try:
        doc = Document(path)
        images = []
        max_images = config.get("max_images", 50)

        for rel in doc.part.rels.values():
            if len(images) >= max_images:
                break
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    img = Image.open(io.BytesIO(image_data))
                    images.append(img)
                except Exception:
                    continue

        if not images:
            return ParseResult.failed(
                "No extractable images found in DOCX", modality="image"
            )

        return ParseResult(
            modality="image",
            output=images,
            metadata={"image_count": len(images), "source_format": "docx"},
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="image")


registry.register([".docx", ".doc"], "image", parse_docx_image)


# ===================================================================
# PPTX
# ===================================================================

def parse_pptx_text(path: str, config: dict, services: dict = None) -> ParseResult:
    """Extract text from a PowerPoint. Detects embedded images."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.debug("python-pptx not installed")
        return ParseResult.failed("python-pptx not installed", modality="text")

    try:
        t0 = time.time()
        limit = _max_chars(config)
        prs = Presentation(path)

        text_runs = []
        current_len = 0
        image_count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text)
                    current_len += len(shape.text)
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image_count += 1
            if current_len > limit:
                break

        content = _clean_text("\n".join(text_runs)[:limit])

        also_contains = []
        if image_count > 0:
            also_contains.append("image")

        logger.debug(
            f"PPTX parsed: {Path(path).name} — {len(prs.slides)} slides, "
            f"{len(content)} chars in {time.time() - t0:.2f}s"
        )
        return ParseResult(
            modality="text",
            output=content,
            metadata={
                "char_count": len(content),
                "slide_count": len(prs.slides),
                "image_count": image_count,
                "has_images": image_count > 0,
            },
            also_contains=also_contains,
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="text")


registry.register(".pptx", "text", parse_pptx_text)


def parse_pptx_image(path: str, config: dict, services: dict = None) -> ParseResult:
    """Extract embedded images from a PPTX as PIL.Image objects."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image
        import io
    except ImportError as e:
        logger.debug(f"Missing dependency: {e}")
        return ParseResult.failed(f"Missing dependency: {e}", modality="image")

    try:
        prs = Presentation(path)
        images = []
        max_images = config.get("max_images", 50)

        for slide in prs.slides:
            if len(images) >= max_images:
                break
            for shape in slide.shapes:
                if len(images) >= max_images:
                    break
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_data = shape.image.blob
                        img = Image.open(io.BytesIO(image_data))
                        images.append(img)
                    except Exception:
                        continue

        if not images:
            return ParseResult.failed(
                "No extractable images found in PPTX", modality="image"
            )

        return ParseResult(
            modality="image",
            output=images,
            metadata={"image_count": len(images), "source_format": "pptx"},
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="image")


registry.register(".pptx", "image", parse_pptx_image)

# ===================================================================
# GDOC - if GoogleDriveService is unloaded, returns False
# ===================================================================


def parse_gdoc(path: str, config: dict, services: dict = None) -> ParseResult:
    """
    Parse a .gdoc file (JSON shortcut) and fetch content from Google Drive.
    Requires the "google_drive" service to be loaded.
    """
    import json

    drive_svc = services.get("google_drive") if services else None

    if drive_svc is None or not getattr(drive_svc, "loaded", False):
        return ParseResult.failed(
            "Drive service not loaded — retry after loading",
            modality="text",
        )

    try:
        with open(path, "r", encoding="utf-8") as f:
            gdoc_data = json.load(f)

        doc_id = gdoc_data.get("doc_id")
        if not doc_id:
            return ParseResult.failed("No doc_id found in .gdoc file", modality="text")

        # Service handles the API call and thread safety internally
        content = drive_svc.download_text(doc_id)

        if content is None:
            return ParseResult.failed("Failed to download document", modality="text")

        limit = _max_chars(config)
        content = _clean_text(content[:limit])

        return ParseResult(
            modality="text",
            output=content,
            metadata={
                "char_count": len(content),
                "source": "google_drive",
                "doc_id": doc_id,
            },
        )
    except Exception as e:
        logger.error(f"Failed to parse gdoc {Path(path).name}: {e}")
        return ParseResult.failed(str(e), modality="text")

registry.register(".gdoc", "text", parse_gdoc)
